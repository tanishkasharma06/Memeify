# app.py
import streamlit as st
import pdfplumber
from pptx import Presentation
from openai import OpenAI

# Initialize OpenAI client (replace with your API key)
client = OpenAI(api_key="YOUR_API_KEY")

# ----------- File Parsing Helpers -----------
def extract_text_from_pdf(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_text_from_pptx(file):
    text = ""
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# ----------- AI Meme Notes Generator -----------
def generate_meme_notes(text):
    prompt = f"""
    You are a fun but smart study assistant.
    Take these class notes:
    {text}

    Step 1: Summarize them into short, clear bullet-point study notes.
    Step 2: Rewrite those notes in meme-style one-liners
            (Gen Z humor, witty, still educational).
    Give both versions separately.
    """
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message["content"]

# ----------- Streamlit App UI -----------
st.set_page_config(page_title="Meme Notes Generator", page_icon="😂")
st.title("📚 Meme Notes Generator")
st.caption("Upload boring notes → Get study notes & meme-fied notes instantly.")

uploaded_file = st.file_uploader("Upload your notes (PDF/PPT)", type=["pdf","pptx"])

if uploaded_file:
    # Extract text
    if uploaded_file.type == "application/pdf":
        text = extract_text_from_pdf(uploaded_file)
    else:
        text = extract_text_from_pptx(uploaded_file)

    if text.strip() == "":
        st.error("⚠️ Could not extract text from file. Try another one.")
    else:
        st.success("✅ File uploaded successfully! Generating meme notes...")
        with st.spinner("AI cooking your memes... 🍳"):
            output = generate_meme_notes(text)
        st.markdown("### ✨ Output")
        st.write(output)

        # Option to download notes
        st.download_button("Download Notes", output, file_name="meme_notes.txt")
